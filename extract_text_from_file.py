import csv
import os
import json
from docx import Document
from pptx import Presentation
import PyPDF2

from extract_ocr_from_docx import extract_ocr_from_docx

# 1. 에러 처리 데코레이터 정의
def handle_extraction_errors(func):
    """텍스트 추출 함수의 예외를 처리하는 데코레이터"""
    def wrapper(filepath):
        try:
            return func(filepath)
        except Exception as e:
            # 함수의 이름에서 파일 형식(예: 'csv')을 가져옴
            file_type = func.__name__.split('_')[1].upper()
            return f"Error reading {file_type}: {str(e)}"
    return wrapper

# 2. 각 추출 함수에 데코레이터 적용
@handle_extraction_errors
def extract_csv_text(filepath):
    """CSV 파일의 텍스트를 추출"""
    with open(filepath, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        return "\n".join([", ".join(row) for row in reader])

@handle_extraction_errors
def extract_docx_text(filepath):
    """DOCX 파일의 본문 텍스트 + 내장 이미지 OCR 결과를 함께 추출"""
    doc = Document(filepath)
    body = "\n".join(para.text for para in doc.paragraphs)

    ocr_texts = extract_ocr_from_docx(filepath)
    if ocr_texts:
        body += "\n\n[Images OCR]\n" + "\n".join(ocr_texts)
    return body

@handle_extraction_errors
def extract_ipynb_text(filepath):
    """Jupyter Notebook(.ipynb) 파일의 텍스트를 추출"""
    with open(filepath, 'r', encoding='utf-8') as f:
        nb = json.load(f)
    text = []
    for cell in nb.get('cells', []):
        if cell.get('cell_type') in ['markdown', 'code']:
            text.append("".join(cell.get('source', [])))
    return "\n".join(text)

@handle_extraction_errors
def extract_md_text(filepath):
    """Markdown 파일의 텍스트를 추출"""
    with open(filepath, 'r', encoding='utf-8') as f:
        return f.read()

@handle_extraction_errors
def extract_pdf_text(filepath):
    """PDF 파일의 텍스트를 추출"""
    with open(filepath, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        text = []
        for page in reader.pages:
            if page_text := page.extract_text():
                text.append(page_text)
    return "\n".join(text)

@handle_extraction_errors
def extract_ppt_text(filepath):
    """PPT/PPTX/PPTM 파일의 텍스트를 추출"""
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# 3. 딕셔너리를 사용하여 확장자와 함수를 매핑
EXTRACTORS = {
    '.csv': extract_csv_text,
    '.docx': extract_docx_text,
    '.ipynb': extract_ipynb_text,
    '.md': extract_md_text,
    '.pdf': extract_pdf_text,
    '.ppt': extract_ppt_text,
    '.pptx': extract_ppt_text,
    '.pptm': extract_ppt_text,
}

def _save_text(text, output_path, source):
    """추출 결과를 .md(기계적 구조화) 또는 .txt로 저장."""
    if output_path.lower().endswith('.md'):
        title = os.path.splitext(os.path.basename(source))[0]
        body = text.replace('[Images OCR]', '## Images OCR\n')
        content = f"# {title}\n\n{body}\n"
    else:
        content = text
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(content)


def extract_text_from_file(filepath, output_path=None):
    """파일 경로를 받아 적절한 텍스트 추출 함수를 호출.

    output_path가 주어지면 결과를 해당 경로에 저장합니다.
    확장자가 .md이면 파일명을 H1로, [Images OCR] 섹션을 H2로 감싼 마크다운으로 저장합니다.
    """
    ext = os.path.splitext(filepath)[1].lower()
    extractor_func = EXTRACTORS.get(ext)
    if extractor_func:
        text = extractor_func(filepath)
    else:
        text = "Unsupported file type"

    if output_path:
        _save_text(text, output_path, source=filepath)

    return text


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="다양한 형식의 파일(.docx, .pdf, .pptx, .csv, .md, .ipynb 등)에서 텍스트를 추출합니다."
    )
    parser.add_argument("filepath", help="입력 파일 경로")
    parser.add_argument("-o", "--output", help="결과 저장 경로 (.md 또는 .txt). 생략하면 stdout에 출력")
    args = parser.parse_args()

    result = extract_text_from_file(args.filepath, output_path=args.output)
    if not args.output:
        print(result)
